from bs4 import BeautifulSoup
import requests
import openai
import scrapy
from scrapy.crawler import CrawlerProcess
from openpyxl import Workbook
from docx import Document
from pptx import Presentation
from nltk.tokenize import word_tokenize
from nltk.tag import pos_tag
from openai import Completion


# Initialize OpenAI API
openai.api_key = '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'


def user_fill_name():
    # This function will take input from the user and proceed with further processing
    while True:
        user = input("Please enter the name of the file: ")
        if user:
            print("Input received:", user)
            confirmation = input("Is this correct? (y for yes, n for no): ")
            if confirmation.lower() == "y":
                # If the user confirms, proceed further
                # You can add your logic here
                print("Proceeding further...")
                break  # Exit the loop and continue further processing
            elif confirmation.lower() == "n":
                # If the user indicates the input is incorrect, repeat the loop
                print("Please try again...")
        else:
            print("Please enter a valid name.")

    return user

def get_fill_type():
    # Function to get the file type from the user
    while True:
        fill_type_input = input("Please enter the file type (excel, word, powerpoint): ").lower()
        if fill_type_input in ['excel', 'word', 'powerpoint']:
            return fill_type_input
        else:
            print("Invalid file type. Please enter 'excel', 'word', or 'powerpoint'.")

def create_file(user_fill_name, fill_type, urls):
    # Function to create the specified type of file
    fill_type = fill_type.lower()

    try:
        if fill_type == 'excel':
            wb = Workbook()
            ws = wb.active
            data = fetch_data_from_internet(urls[0])
            wb.save(f'{user_fill_name}.xlsx')
            ws.cell(row=1, column=1, value=data) # type: ignore
            print("Excel file created successfully.")

        elif fill_type == 'word':
            doc = Document()
            text = fetch_data_from_internet(urls[1])
            doc.add_paragraph(text)
            doc.save(f'{user_fill_name}.docx')
            print("Word document created successfully.")

        elif fill_type == 'powerpoint':
            prs = Presentation()
            text = fetch_data_from_internet(urls[2])
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = text
            prs.save(f'{user_fill_name}.pptx')
            print("PowerPoint presentation created successfully.")

        else:
            print("Invalid file type. Please choose 'excel', 'word', or 'powerpoint'.")

    except Exception as e:
        print(f"An error occurred: {e}")

def fill_type_handler(urls):
    # Function to handle the fill type and file creation
    print("Select the file type:")
    fill_type = get_fill_type()
    if fill_type:
        create_file(user_fill_name(), fill_type, urls)

def get_user_input():
    user_input = input("Please enter your needs: ")
    return user_input

def process_user_input(user_input):
    # Use OpenAI's GPT model to process user input
    response = openai.Completion.create(
        engine="text-davinci-003",  # You can choose the desired model
        prompt=user_input,
        max_tokens=50  # Adjust the max_tokens based on your requirements
    )

    # Extract entities and intent from the generated text
    processed_text = response.choices[0].text.strip()
    tokens = word_tokenize(processed_text)
    tagged_tokens = pos_tag(tokens)
    entities = [word for word, tag in tagged_tokens if tag.startswith('NN')]
    intent = [word for word, tag in tagged_tokens if tag.startswith('VB')]

    return entities, intent

def fetch_data_from_internet(urls):
    try:
        # Send a GET request to the URL
        response = requests.get(urls)
        
        # Check if the request was successful
        if response.status_code == 200:
            # Parse the HTML content using BeautifulSoup
            soup = BeautifulSoup(response.text, 'html.parser')
            
            # Extract the data you need from the parsed HTML
            # For example, extract text from paragraphs
            paragraphs = [p.text for p in soup.find_all('p')]
            
            return paragraphs
        else:
            print("Failed to fetch data. Status code:", response.status_code)
            return None
    except Exception as e:
        print("An error occurred:", e)
        return None

class MySpider(scrapy.Spider):
    name = 'myspider'
    
    def __init__(self, user_input, **kwargs):
        self.user_input = user_input
        super().__init__(**kwargs)
    
    def generate_urls(self):   
        # Define the URLs or search queries based on the user input
        # For example, if the user input is a search query, construct the search URL
        urls = [
            f'https://www.google.com/search?q={self.user_input}',
            f'https://download-directory.github.io/?query={self.user_input}',
            f'https://en.wikipedia.org/?q={self.user_input}',
            f'https://www.amazon.com?q={self.user_input}',
            f'https://chat.openai.com/',
            # Add more URLs here as needed
        ]
        
        # Iterate over the URLs and yield requests for each
        for url in urls:
            yield scrapy.Request(url=url, callback=self.parse)
        return urls 

    def parse(self, response):
        # Initialize variables with default values
        soup = None
        title = None

        # Use try-except block to handle exceptions
        try:
            # Use BeautifulSoup to parse the HTML content
            soup = BeautifulSoup(response.text, 'html.parser')

            # Extract data from the parsed HTML
            # For example, extract the title if it exists
            if soup.title:
                title = soup.title.string

        except Exception as e:
            # Handle any exceptions that may occur during parsing
            print("An error occurred while parsing:", e)

            # Extract URLs from the metadata
            urls = response.meta.get('urls', [])

            # Call create_file function with the extracted URLs
            create_file(user_fill_name, get_fill_type, urls)   


        # Generate output with OpenAI
        output = generate_output({'title': title}, self.user_input)
        yield {'output': output}



def generate_output(data, user_input):
    # Use OpenAI to generate output based on the scraped data and user input
    # For example, generate a summary or a report based on the data and user input
    prompt = f"Generate content based on the scraped data and user input:\n\nScraped Title: {data['title']}\nScraped Paragraphs: {data['paragraphs']}\n\nUser Input: {user_input}\n\nOutput:"
    response = openai.Completion.create( # type: ignore
        engine="text-davinci-003",
        prompt=prompt,
        max_tokens=200
    )
    return response.choices[0].text.strip()

def scrape_data(user_input):
    # Create a CrawlerProcess
    process = CrawlerProcess()
    # Add the spider to the process
    process.crawl(MySpider, user_input=user_input)
    # Start the crawling process
    process.start()    

if __name__ == "__main__":
    # Get user input
    user_input = get_user_input()
    
    # Process user input
    entities, intent = process_user_input(user_input)
    
    # Scrape data based on user input
    scrape_data(user_input)
    
    # Ask user for file name
    file_name = user_fill_name()
    
    # Ask user for file type and create the file
   